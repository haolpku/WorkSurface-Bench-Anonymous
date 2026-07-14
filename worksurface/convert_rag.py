"""RAG surface: extract canonical text from workspace documents.

Every document-type input file (md, txt, pdf, docx, pptx, doc/ppt, html) is
converted to a single canonical UTF-8 text/markdown file under a profile's
``kb_docs/`` dir. Tabular files are handled by the Table surface, not here —
but a tabular sheet that fails the Table coverage gate can be *demoted* to
RAG (rendered as a markdown table); that demotion is decided in
convert_tables, which calls :func:`csv_to_markdown` from here.

Canonicalization is deliberate (a paper design point, README "Canonical
surfaces"): routing is evaluated separately from OCR / raw parsing, so we
freeze the extracted text once and score against it.
"""

from __future__ import annotations

import csv
import io
import os

from .common import Task, strip_hash_prefix

# Extensions this surface knows how to turn into text.
RAG_EXTS = {".md", ".txt", ".pdf", ".docx", ".pptx", ".doc", ".ppt", ".html", ".java", ".py"}


def _read_text(path: str) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            with open(path, encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, ValueError):
            continue
    with open(path, "rb") as f:
        return f.read().decode("utf-8", errors="replace")


def _pdf_to_text(path: str) -> str:
    try:
        from pypdf import PdfReader

        reader = PdfReader(path)
        parts = []
        for i, page in enumerate(reader.pages):
            txt = page.extract_text() or ""
            if txt.strip():
                parts.append(f"\n\n<!-- page {i + 1} -->\n{txt}")
        return "".join(parts).strip()
    except Exception as e:  # noqa: BLE001 - want a placeholder, not a crash
        return f"[PDF text extraction failed: {type(e).__name__}: {e}]"


def _docx_to_text(path: str) -> str:
    try:
        import docx

        doc = docx.Document(path)
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                lines.append(" | ".join(cells))
        return "\n".join(lines).strip()
    except Exception as e:  # noqa: BLE001
        return f"[DOCX text extraction failed: {type(e).__name__}: {e}]"


def _pptx_to_text(path: str) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
            if texts:
                parts.append(f"\n\n<!-- slide {i + 1} -->\n" + "\n".join(texts))
        return "".join(parts).strip()
    except Exception as e:  # noqa: BLE001
        return f"[PPTX text extraction failed: {type(e).__name__}: {e}]"


def csv_to_markdown(path: str, max_rows: int = 200) -> str:
    """Render a CSV as a GitHub-flavored markdown table (RAG demotion path)."""
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        rows = list(csv.reader(f))
    if not rows:
        return ""
    header, body = rows[0], rows[1:]
    out = ["| " + " | ".join(header) + " |",
           "| " + " | ".join("---" for _ in header) + " |"]
    for r in body[:max_rows]:
        out.append("| " + " | ".join(r) + " |")
    if len(body) > max_rows:
        out.append(f"\n<!-- {len(body) - max_rows} more rows omitted -->")
    return "\n".join(out)


def extract_text(abspath: str, ext: str) -> str:
    if ext in (".md", ".txt", ".html", ".java", ".py"):
        return _read_text(abspath)
    if ext == ".pdf":
        return _pdf_to_text(abspath)
    if ext in (".docx",):
        return _docx_to_text(abspath)
    if ext in (".pptx",):
        return _pptx_to_text(abspath)
    if ext in (".doc", ".ppt"):
        # Legacy binary Office — no pure-python extractor bundled; emit a stub
        # so the doc still exists as a routable (if empty) KB entry.
        return f"[legacy binary format {ext}; text not extracted]"
    return _read_text(abspath)


def build_kb_docs(
    profile_tasks: list[Task],
    profile_dir: str,
    demote_to_rag: set[tuple[str, str]] | None = None,
) -> dict[str, dict]:
    """Write canonical text for every doc file into ``{profile_dir}/kb_docs/``.

    Returns a registry: canonical_doc_name -> {source_task, source_file,
    ext, chars, from_table}. ``demote_to_rag`` is a set of (task_id,
    filename) tabular sheets that failed the Table gate and should be
    rendered as markdown tables here instead.
    """
    kb_dir = os.path.join(profile_dir, "kb_docs")
    os.makedirs(kb_dir, exist_ok=True)
    demote_to_rag = demote_to_rag or set()
    registry: dict[str, dict] = {}

    for task in profile_tasks:
        for entry in task.manifest():
            if not entry["exists"]:
                continue
            ext = entry["ext"]
            clean = strip_hash_prefix(entry["filename"])
            is_demoted = (task.task_id, entry["filename"]) in demote_to_rag

            if ext in RAG_EXTS:
                text = extract_text(entry["abspath"], ext)
            elif is_demoted and ext == ".csv":
                text = csv_to_markdown(entry["abspath"])
            else:
                continue  # tabular file handled by Table surface

            # Namespace by task id to avoid collisions across tasks that reuse
            # generic names like "report.md".
            doc_name = f"t{task.task_id}__{os.path.splitext(clean)[0]}.md"
            out_path = os.path.join(kb_dir, doc_name)
            header = (
                f"<!-- source_task: {task.task_id} | "
                f"source_file: {clean} | surface: rag -->\n\n"
            )
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(header + (text or ""))
            registry[doc_name] = {
                "source_task": task.task_id,
                "source_file": clean,
                "ext": ext,
                "chars": len(text or ""),
                "from_table": is_demoted,
            }
    return registry
